from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_hospital_management_presentation():
    # Create a new presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content
    section_header_layout = prs.slide_layouts[2]  # Section header
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Hospital Management System"
    subtitle.text = "A Comprehensive Web-Based Solution\n\nDeveloped using Django Framework\n\nTeam Members: [Your Names]\nDate: [Current Date]"
    
    # Slide 2: Introduction
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "1. Introduction & Objectives"
    content_text = """• Project Overview:
  - Web-based Hospital Management System
  - Built with Django (Python Framework)
  - SQLite Database backend
  
• Main Objectives:
  - Streamline hospital operations
  - Manage doctor and patient records
  - Schedule and track appointments
  - Provide secure admin interface
  - Generate reports and analytics
  
• Key Features:
  - User authentication and authorization
  - CRUD operations for all entities
  - Real-time dashboard with statistics
  - Responsive web interface"""
    
    content.text = content_text
    
    # Slide 3: SRS - Functional Requirements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "2. System Requirements Specification"
    content_text = """Functional Requirements:
• User Management:
  - Admin login/logout functionality
  - Secure authentication system
  
• Doctor Management:
  - Add new doctors with details
  - View all doctor records
  - Delete doctor records
  - Store: name, specialization, contact
  
• Patient Management:
  - Add new patients with details
  - View all patient records
  - Store: name, gender, mobile, address
  
• Appointment Management:
  - Schedule appointments
  - View all appointments
  - Delete appointments
  - Link patients with doctors

Non-Functional Requirements:
• Performance: Fast response times
• Security: Admin-only access
• Usability: Intuitive interface
• Reliability: Data consistency"""
    
    content.text = content_text
    
    # Slide 4: Process Logic
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "3. Process Logic"
    content_text = """System Workflow:

1. Authentication Process:
   - Admin enters credentials
   - System validates user
   - Redirects to dashboard if valid

2. Doctor Management Process:
   - Admin adds doctor details
   - System validates input
   - Stores in database
   - Updates doctor list

3. Patient Management Process:
   - Admin adds patient information
   - System validates data
   - Creates patient record
   - Updates patient directory

4. Appointment Scheduling:
   - Select patient from list
   - Choose available doctor
   - Set date and time
   - Create appointment record

5. Data Flow:
   - User Input → Validation → Database → Response
   - Real-time dashboard updates
   - Error handling and feedback"""
    
    content.text = content_text
    
    # Slide 5: Gantt Chart
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "4. Project Timeline (3 Months)"
    content_text = """Gantt Chart - 12 Weeks Duration:

Week 1-2: Project Planning & Setup
• Requirements gathering
• Technology stack selection
• Development environment setup

Week 3-4: Database Design
• ER diagram creation
• Database schema design
• Django models implementation

Week 5-6: Core Functionality
• User authentication system
• Admin interface development
• Basic CRUD operations

Week 7-8: Doctor Management
• Doctor registration system
• Doctor listing and search
• Doctor profile management

Week 9-10: Patient Management
• Patient registration system
• Patient records management
• Patient search functionality

Week 11-12: Appointment System
• Appointment scheduling
• Calendar integration
• Appointment tracking

Week 13-14: Testing & Documentation
• Unit testing
• Integration testing
• User documentation

Week 15-16: Deployment & Finalization
• Production deployment
• Performance optimization
• Final testing and bug fixes"""
    
    content.text = content_text
    
    # Slide 6: Data Dictionary
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "5. Data Dictionary"
    content_text = """Database Schema:

DOCTOR Table:
• id (Primary Key): Auto-increment integer
• name: VARCHAR(100) - Doctor's full name
• specialization: VARCHAR(100) - Medical specialty
• contact_number: VARCHAR(20) - Phone number

PATIENT Table:
• id (Primary Key): Auto-increment integer
• name: VARCHAR(100) - Patient's full name
• gender: VARCHAR(10) - Male/Female/Other
• mobile_number: INTEGER - Contact number
• address: VARCHAR(200) - Patient's address

APPOINTMENT Table:
• id (Primary Key): Auto-increment integer
• patient_id (Foreign Key): References Patient.id
• doctor_id (Foreign Key): References Doctor.id
• appointment_date: DATE - Scheduled date
• appointment_time: TIME - Scheduled time

USER Table (Django Auth):
• id, username, password, email, is_staff, etc."""
    
    content.text = content_text
    
    # Slide 7: DFD
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "6. Data Flow Diagram (DFD)"
    content_text = """Level 0 DFD - Context Diagram:

[Admin User] → [Hospital Management System] → [Database]

Level 1 DFD - Main Processes:

1. Authentication Process:
   Admin → Login Form → Validation → Dashboard

2. Doctor Management:
   Admin → Doctor Forms → CRUD Operations → Doctor DB

3. Patient Management:
   Admin → Patient Forms → CRUD Operations → Patient DB

4. Appointment Management:
   Admin → Appointment Forms → Scheduling → Appointment DB

5. Reporting:
   Admin → Dashboard → Statistics → Reports

Data Stores:
• Doctor Database
• Patient Database  
• Appointment Database
• User Authentication DB

External Entities:
• Admin Users
• System Reports"""
    
    content.text = content_text
    
    # Slide 8: ER Diagram
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "7. Entity Relationship Diagram"
    content_text = """ER Diagram Structure:

┌─────────────┐         ┌─────────────┐         ┌─────────────┐
│   DOCTOR    │         │ APPOINTMENT │         │   PATIENT   │
├─────────────┤         ├─────────────┤         ├─────────────┤
│ PK: id      │         │ PK: id      │         │ PK: id      │
│ name        │         │ FK: patient │◄────────┤ name        │
│ specialization│       │ FK: doctor  │         │ gender      │
│ contact_number│       │ date        │         │ mobile_number│
└─────────────┘         │ time        │         │ address     │
         ▲              └─────────────┘         └─────────────┘
         │                       ▲
         └───────────────────────┘

Relationships:
• Doctor (1) ←→ (M) Appointment
• Patient (1) ←→ (M) Appointment
• One doctor can have many appointments
• One patient can have many appointments
• Each appointment belongs to one doctor and one patient

Cardinality:
• Doctor to Appointment: 1:N
• Patient to Appointment: 1:N
• Appointment to Doctor: N:1
• Appointment to Patient: N:1"""
    
    content.text = content_text
    
    # Slide 9: Interface Design
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "8. User Interface Design"
    content_text = """Interface Components:

1. Login Interface:
   - Username and password fields
   - Secure authentication form
   - Error message display

2. Admin Dashboard:
   - Statistics overview (doctors, patients, appointments)
   - Quick access navigation menu
   - Real-time counters

3. Doctor Management Interface:
   - Add Doctor Form: name, specialization, contact
   - View Doctors: table with all doctor records
   - Delete functionality with confirmation

4. Patient Management Interface:
   - Add Patient Form: name, gender, mobile, address
   - View Patients: comprehensive patient list
   - Search and filter capabilities

5. Appointment Management Interface:
   - Schedule Appointment: patient/doctor selection, date/time
   - View Appointments: detailed appointment list
   - Calendar view integration

6. Navigation:
   - Responsive navigation bar
   - Breadcrumb navigation
   - Logout functionality

Design Principles:
• Clean and intuitive interface
• Consistent color scheme
• Mobile-responsive design
• User-friendly forms and validation"""
    
    content.text = content_text
    
    # Slide 10: Expected Reports
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "9. Expected Report Generation"
    content_text = """System Reports:

1. Dashboard Statistics:
   • Total number of doctors
   • Total number of patients
   • Total appointments scheduled
   • Real-time system overview

2. Doctor Reports:
   • Complete doctor directory
   • Doctor specialization analysis
   • Doctor availability status
   • Doctor performance metrics

3. Patient Reports:
   • Patient demographic analysis
   • Patient registration trends
   • Patient visit frequency
   • Patient contact directory

4. Appointment Reports:
   • Daily appointment schedule
   • Weekly/monthly appointment trends
   • Doctor appointment load
   • Appointment completion status

5. Administrative Reports:
   • System usage statistics
   • User activity logs
   • Data backup reports
   • System performance metrics

6. Custom Reports:
   • Date range specific reports
   • Specialization-wise analysis
   • Patient-doctor relationship reports
   • Appointment conflict reports"""
    
    content.text = content_text
    
    # Slide 11: References
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "10. References & Bibliography"
    content_text = """Technical References:

1. Django Documentation:
   • https://docs.djangoproject.com/
   • Django 5.2.3 Official Documentation
   • Django ORM and Models Guide

2. Python Documentation:
   • https://docs.python.org/
   • Python 3.x Programming Guide
   • PEP 8 Style Guidelines

3. Web Development:
   • HTML5 and CSS3 Standards
   • Bootstrap Framework Documentation
   • JavaScript ES6+ Features

4. Database Design:
   • SQLite Documentation
   • Database Normalization Principles
   • ER Diagram Design Guidelines

5. Hospital Management Systems:
   • Healthcare Information Systems
   • Patient Data Management Standards
   • Medical Appointment Systems

6. Security References:
   • Django Security Best Practices
   • OWASP Security Guidelines
   • Authentication and Authorization

7. Testing and Deployment:
   • Django Testing Framework
   • Web Application Testing
   • Production Deployment Guidelines"""
    
    content.text = content_text
    
    # Slide 12: Future Scope
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "11. Future Scope & Enhancements"
    content_text = """Potential Enhancements:

1. Advanced Features:
   • Multi-user roles (Doctors, Nurses, Receptionists)
   • Patient medical history tracking
   • Prescription management system
   • Billing and payment integration
   • Inventory management for medical supplies

2. Technology Upgrades:
   • Mobile application development
   • API development for third-party integration
   • Cloud deployment and scaling
   • Real-time notifications system
   • Advanced analytics and reporting

3. Healthcare Integration:
   • Electronic Health Records (EHR)
   • Laboratory management system
   • Radiology and imaging integration
   • Pharmacy management system
   • Insurance claim processing

4. Security Enhancements:
   • Two-factor authentication
   • Data encryption at rest
   • HIPAA compliance features
   • Audit trail and logging
   • Backup and disaster recovery

5. User Experience:
   • Advanced search and filtering
   • Calendar integration
   • Email and SMS notifications
   • Multi-language support
   • Accessibility improvements

6. Analytics and AI:
   • Predictive analytics for patient flow
   • AI-powered appointment scheduling
   • Disease pattern analysis
   • Resource optimization
   • Performance monitoring dashboards"""
    
    content.text = content_text
    
    # Slide 13: Conclusion
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "12. Conclusion"
    content_text = """Project Summary:

✅ Achievements:
• Successfully developed a functional Hospital Management System
• Implemented secure user authentication
• Created comprehensive CRUD operations
• Built responsive web interface
• Established proper database relationships

✅ Key Benefits:
• Streamlined hospital operations
• Improved data management
• Enhanced appointment scheduling
• Better resource utilization
• Increased administrative efficiency

✅ Technical Accomplishments:
• Django framework implementation
• SQLite database integration
• MVC architecture pattern
• RESTful API design principles
• Responsive web design

✅ Learning Outcomes:
• Full-stack web development
• Database design and management
• User interface design
• System architecture planning
• Project management skills

The Hospital Management System provides a solid foundation for healthcare facility management and can be extended with additional features as per future requirements."""
    
    content.text = content_text
    
    # Save the presentation
    prs.save('Hospital_Management_System_Presentation.pptx')
    print("Presentation created successfully: Hospital_Management_System_Presentation.pptx")

if __name__ == "__main__":
    create_hospital_management_presentation() 